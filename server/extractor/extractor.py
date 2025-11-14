from flask import Flask, request, jsonify
import io, re

app = Flask(__name__)

def extract_from_pptx(file_bytes):
    from pptx import Presentation
    from PIL import Image
    try:
        import pytesseract
    except Exception:
        pytesseract = None

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        entry = {"slide": i, "title": "", "bullets": [], "notes": "", "images_text": []}
        for idx, shape in enumerate(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if not entry["title"] and idx == 0:
                    entry["title"] = text
                else:
                    parts = [p.strip() for p in re.split(r"\n|•|- ", text) if p.strip()]
                    entry["bullets"].extend(parts)
            if getattr(shape, "shape_type", None) == 13 and pytesseract:  # PICTURE
                try:
                    from PIL import Image
                    img = Image.open(io.BytesIO(shape.image.blob))
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text and ocr_text.strip():
                        entry["images_text"].append(ocr_text.strip())
                except Exception:
                    pass
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide and slide.notes_slide.notes_text_frame:
            entry["notes"] = (slide.notes_slide.notes_text_frame.text or "").strip()
        slides.append(entry)
    return slides

def extract_from_pdf(file_bytes):
    import fitz  # PyMuPDF
    slides = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        bullets = [t.strip() for t in text.split("\n") if t.strip()]
        slides.append({"slide": i, "title": "", "bullets": bullets, "notes": "", "images_text": []})
    return slides

@app.post("/extract")
def extract():
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "file missing"}), 400
    buf = f.read()
    name = (f.filename or "").lower()

    if name.endswith(".pptx"):
        slides = extract_from_pptx(buf)
    elif name.endswith(".pdf"):
        slides = extract_from_pdf(buf)
    else:
        return jsonify({"error": "unsupported"}), 400

    return jsonify({"slides": slides, "count": len(slides)})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5001)
